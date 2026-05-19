from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SLIDES = ROOT / "slides"
FIGURES = ROOT / "outputs" / "figures"
OUT = SLIDES / "midterm_report.pptx"

W = Inches(13.333)
H = Inches(7.5)
NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(14, 165, 233)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_text(slide, text, x, y, w, h, size=18, color=SLATE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.clear()
    for i, line in enumerate(str(text).split("\n")):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        if align:
            p.alignment = align
    return box


def add_title(slide, title, kicker=None):
    if kicker:
        add_text(slide, kicker, 0.65, 0.28, 4.8, 0.25, 9, BLUE, True)
    add_text(slide, title, 0.62, 0.48, 8.6, 0.55, 24, NAVY, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(1.08), Inches(1.0), Inches(0.04))
    set_fill(line, BLUE)


def add_footer(slide, n):
    add_text(slide, f"AI Data Foundation 中期汇报 · {n}", 10.6, 7.08, 2.1, 0.2, 8, MUTED, False, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent=BLUE, title_size=14, body_size=12):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.03), Inches(y + 0.04), Inches(w), Inches(h))
    set_fill(shadow, RGBColor(226, 232, 240))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(card, WHITE)
    card.line.color.rgb = BORDER
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    set_fill(bar, accent)
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.35, 0.35, title_size, NAVY, True)
    add_text(slide, body, x + 0.22, y + 0.62, w - 0.35, h - 0.75, body_size, SLATE)


def add_bullets(slide, bullets, x, y, w, h, size=16, color=SLATE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def add_metric(slide, x, y, w, h, value, label, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, WHITE)
    shape.line.color.rgb = RGBColor(203, 213, 225)
    add_text(slide, value, x + 0.12, y + 0.16, w - 0.24, 0.38, 22, accent, True, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.12, y + 0.62, w - 0.24, 0.3, 9, MUTED, False, PP_ALIGN.CENTER)


def add_image_panel(slide, image_name, x=6.65, y=1.35, w=5.95, h=4.4):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(panel, WHITE)
    panel.line.color.rgb = BORDER
    path = FIGURES / image_name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x + 0.22), Inches(y + 0.28), width=Inches(w - 0.44))


def add_section(slide, title, subtitle, idx):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    set_fill(bg, NAVY)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.0), Inches(4.5), Inches(4.5))
    set_fill(circle, RGBColor(30, 64, 175))
    add_text(slide, f"0{idx}", 0.85, 1.55, 1.0, 0.55, 24, CYAN, True)
    add_text(slide, title, 0.85, 2.2, 8.5, 0.8, 34, WHITE, True)
    add_text(slide, subtitle, 0.9, 3.25, 8.8, 0.6, 17, RGBColor(203, 213, 225))


def add_arrow(slide, x, y, w, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(x), Inches(y), Inches(w), Inches(0.75))
    set_fill(shape, color)
    add_text(slide, text, x + 0.08, y + 0.22, w - 0.18, 0.22, 11, WHITE, True, PP_ALIGN.CENTER)


def add_table_like(slide, rows, x, y, widths, row_h=0.42, header=True):
    for r, row in enumerate(rows):
        left = x
        for c, text in enumerate(row):
            fill = NAVY if r == 0 and header else (LIGHT if r % 2 == 0 else WHITE)
            color = WHITE if r == 0 and header else SLATE
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(y + r * row_h), Inches(widths[c]), Inches(row_h))
            set_fill(rect, fill)
            rect.line.color.rgb = BORDER
            add_text(slide, text, left + 0.05, y + r * row_h + 0.08, widths[c] - 0.1, row_h - 0.08, 9.5, color, r == 0 and header)
            left += widths[c]


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    set_fill(bg, NAVY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), H)
    set_fill(accent, CYAN)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.7), Inches(-1.2), Inches(5.5), Inches(5.5))
    set_fill(circle, RGBColor(30, 64, 175))
    add_text(slide, "AI DATA FOUNDATION · MIDTERM REPORT", 0.9, 0.9, 6.8, 0.3, 10, CYAN, True)
    add_text(slide, "基于多源数据的个人贷款\n违约风险检测", 0.85, 1.55, 8.6, 1.65, 38, WHITE, True)
    add_text(slide, "从 Lending Club 真实贷款记录出发，融合地区经济与宏观金融变量，构建可解释的违约风险分析框架。", 0.9, 3.45, 7.4, 0.7, 17, RGBColor(203, 213, 225))
    add_metric(slide, 0.9, 5.15, 2.2, 1.05, "226万", "accepted 贷款记录", CYAN)
    add_metric(slide, 3.35, 5.15, 2.2, 1.05, "136.8万", "已完结样本", GREEN)
    add_metric(slide, 5.8, 5.15, 2.2, 1.05, "21.27%", "总体违约率", ORANGE)
    add_text(slide, "日期：2026-05-21\n小组成员：请替换为实际姓名", 0.9, 6.65, 5.2, 0.5, 12, RGBColor(203, 213, 225))


def slide_motivation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "研究动机：为什么这件事值得做", "01 / Motivation")
    add_card(slide, 0.75, 1.45, 3.75, 2.0, "现实价值", "个人贷款违约会直接影响坏账损失、资金定价和借款人信用可得性。", BLUE)
    add_card(slide, 4.8, 1.45, 3.75, 2.0, "传统不足", "只看借款人自身画像，容易忽略地区就业、收入水平和宏观金融周期。", ORANGE)
    add_card(slide, 8.85, 1.45, 3.75, 2.0, "项目主张", "违约风险 = 个体信用特征 + 地区经济环境 + 宏观金融周期。", GREEN)
    add_text(slide, "一句话主线", 0.8, 4.15, 1.7, 0.3, 13, BLUE, True)
    add_text(slide, "同样收入、同样信用等级的人，在不同地区和不同经济周期下，违约风险可能完全不同。", 0.78, 4.55, 11.7, 0.55, 26, NAVY, True)
    add_footer(slide, 2)


def slide_questions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "研究问题与挑战", "02 / Problem")
    add_card(slide, 0.75, 1.35, 3.7, 1.45, "RQ1 数据融合", "如何整合个人信贷、地区经济与宏观金融数据？", BLUE)
    add_card(slide, 4.82, 1.35, 3.7, 1.45, "RQ2 粒度对齐", "如何处理个人、州级、年度/季度级数据之间的口径差异？", PURPLE)
    add_card(slide, 8.88, 1.35, 3.7, 1.45, "RQ3 解释价值", "多源变量是否能提供比单一贷款数据更可靠的风险洞察？", GREEN)
    add_bullets(slide, ["数据规模大：2,260,701 条 accepted 贷款记录，151 个字段。", "标签复杂：Current 等未完结贷款不能直接判断最终违约。", "结果需要可解释：不能只给模型分数，还要说明风险来源。"], 1.0, 3.55, 11.2, 1.7, 18)
    add_footer(slide, 3)


def slide_related(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "相关工作：我们站在哪里", "03 / Related Work")
    add_card(slide, 0.75, 1.35, 3.7, 3.65, "传统信用评分", "FICO、信用历史、收入、负债率\n\n优点：解释性强、金融机构熟悉\n局限：偏重个体特征，难以捕捉外部环境", BLUE)
    add_card(slide, 4.82, 1.35, 3.7, 3.65, "机器学习预测", "Logistic Regression、Random Forest、XGBoost\n\n优点：预测能力强\n局限：容易黑盒化，业务解释不足", PURPLE)
    add_card(slide, 8.88, 1.35, 3.7, 3.65, "多源风险分析", "加入收入、失业率、利率、通胀等外部变量\n\n优点：视角全面\n局限：数据对齐和口径一致性困难", GREEN)
    add_text(slide, "本项目定位：先把真实数据清洗、融合、统计分析和可解释洞察做好，再考虑预测模型。", 1.0, 5.55, 11.3, 0.5, 18, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 4)


def slide_data_status(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "数据集与中期状态", "04 / Data Assets")
    rows = [
        ["数据源", "作用", "中期状态"],
        ["Lending Club", "贷款状态、等级、利率、收入、州、时间", "已下载并完成真实 EDA"],
        ["Home Credit", "补充信贷场景对照", "需先接受 Kaggle competition rules"],
        ["ACS", "州/县收入、贫困率、就业", "下一阶段接入"],
        ["Fed", "利率、通胀、失业率", "下一阶段接入"],
    ]
    add_table_like(slide, rows, 0.85, 1.45, [2.3, 5.0, 4.5], 0.72)
    add_metric(slide, 1.1, 5.45, 2.25, 0.95, "151", "Lending Club 字段数", BLUE)
    add_metric(slide, 3.75, 5.45, 2.25, 0.95, "2.26M", "原始记录", CYAN)
    add_metric(slide, 6.4, 5.45, 2.25, 0.95, "1.37M", "已完结样本", GREEN)
    add_metric(slide, 9.05, 5.45, 2.25, 0.95, "21.27%", "总体违约率", ORANGE)
    add_footer(slide, 5)


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "数据流水线：从原始 CSV 到可解释图表", "05 / Pipeline")
    steps = [("Kaggle\n原始数据", BLUE), ("字段检查\n缺失统计", CYAN), ("状态过滤\n标签定义", GREEN), ("分组聚合\n违约率", ORANGE), ("表格图表\n汇报产出", PURPLE)]
    for i, (text, color) in enumerate(steps):
        add_arrow(slide, 0.75 + i * 2.48, 2.0, 2.05, text, color)
    add_card(slide, 1.0, 4.05, 5.35, 1.55, "关键处理", "只用 Fully Paid / Charged Off / Default 等可判断最终结果的状态计算违约率，避免把 Current 误判为不违约。", BLUE)
    add_card(slide, 6.95, 4.05, 5.35, 1.55, "下一步扩展", "把 ACS 州级经济指标和 Fed 宏观金融指标接入，解释地区差异和时间趋势。", GREEN)
    add_footer(slide, 6)


def slide_progress(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "当前进度：已经从框架推进到真实数据", "06 / Progress")
    add_metric(slide, 0.85, 1.35, 2.2, 1.0, "GitHub", "项目仓库已上传", BLUE)
    add_metric(slide, 3.35, 1.35, 2.2, 1.0, "2.26M", "Lending Club 已读取", CYAN)
    add_metric(slide, 5.85, 1.35, 2.2, 1.0, "1.37M", "已完结记录", GREEN)
    add_metric(slide, 8.35, 1.35, 2.2, 1.0, "5+", "真实数据图表", ORANGE)
    add_bullets(slide, ["完成代码结构、README、数据说明、PPT 生成脚本。", "完成 Lending Club 大文件流式分析，输出真实统计表与图表。", "当前节点：单一真实信贷数据源分析完成，下一步进入 ACS/Fed 多源融合。"], 1.1, 3.05, 10.8, 1.5, 18)
    add_footer(slide, 7)


def slide_finding(prs, number, title, bullets, image, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, f"真实发现 {number}：{title}", f"Finding {number}")
    add_card(slide, 0.85, 1.45, 5.35, 3.8, "核心解读", "\n".join(bullets), accent, 15, 13)
    add_image_panel(slide, image)
    add_footer(slide, 7 + number)


def slide_attempts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "哪些尝试奏效，哪些没有", "12 / What Worked")
    add_card(slide, 0.85, 1.35, 5.65, 4.55, "奏效", "1. Lending Club 大型 CSV 成功流式处理\n2. 贷款状态过滤避免标签误判\n3. 真实数据图表可直接用于中期汇报\n4. GitHub 仓库和运行脚本可复现", GREEN)
    add_card(slide, 6.85, 1.35, 5.65, 4.55, "暂未奏效 / 受阻", "1. Home Credit 需要网页端接受 Kaggle rules\n2. ACS/Fed 尚未真实融合\n3. 当前可视化仍以静态图为主\n4. 多源变量贡献需要下一阶段验证", ORANGE)
    add_footer(slide, 13)


def slide_gantt(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "项目甘特图与当前节点", "13 / Timeline")
    rows = [
        ("开题与问题定义", "4/20–4/30", 0.6, 1.4, GREEN),
        ("项目仓库与数据说明", "5/15–5/18", 3.0, 0.55, GREEN),
        ("Lending Club EDA", "5/18–5/19", 3.55, 0.45, GREEN),
        ("Home Credit / ACS / Fed", "5/20–5/27", 3.95, 1.5, ORANGE),
        ("多源融合与对比", "5/27–6/3", 5.35, 1.6, BLUE),
        ("最终报告与视频", "6/3–6/17", 6.95, 2.9, PURPLE),
    ]
    add_text(slide, "4/20", 3.0, 1.25, 0.6, 0.2, 9, MUTED)
    add_text(slide, "5/19 当前", 5.0, 1.25, 1.0, 0.2, 9, ORANGE, True)
    add_text(slide, "6/17", 9.75, 1.25, 0.6, 0.2, 9, MUTED)
    for i, (name, dates, start, span, color) in enumerate(rows):
        y = 1.75 + i * 0.68
        add_text(slide, name, 0.75, y + 0.08, 2.25, 0.22, 10, NAVY, True)
        add_text(slide, dates, 2.55, y + 0.08, 1.0, 0.22, 9, MUTED)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.65 + start), Inches(y), Inches(span), Inches(0.32))
        set_fill(bar, color)
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.55), Inches(1.55), Inches(0.035), Inches(4.5))
    set_fill(marker, RED)
    add_text(slide, "当前：真实 Lending Club 分析完成，准备进入多源融合", 1.0, 6.25, 11.0, 0.35, 17, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 14)


def slide_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "后续计划：把单源洞察推进到多源解释", "14 / Next Steps")
    rows = [
        ["时间", "任务", "产出"],
        ["5/20–5/22", "Home Credit 授权下载；若失败则作为补充", "数据可用性结论"],
        ["5/22–5/27", "下载 ACS/Fed，统一州代码和年份/季度口径", "外部变量表"],
        ["5/27–6/3", "Lending Club 与 ACS/Fed 州级、时间级融合", "多源融合数据集"],
        ["6/3–6/8", "最终可视化、核心结论、报告初稿", "报告和图表"],
        ["6/8–6/17", "打磨最终 PPT、视频和代码仓库", "最终提交物"],
    ]
    add_table_like(slide, rows, 0.65, 1.35, [1.7, 6.2, 3.6], 0.67)
    add_footer(slide, 15)


def slide_risks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "风险与备选方案", "15 / Risks")
    rows = [
        ["风险", "影响", "备选方案"],
        ["Home Credit 需网页授权", "数据无法按期下载", "先以 Lending Club + ACS + Fed 完成主线"],
        ["ACS/Fed 粒度不一致", "JOIN 质量不稳定", "降到州级、年度/季度级对齐"],
        ["本地处理大文件慢", "影响迭代速度", "流式 CSV 统计 + 抽样分析"],
        ["多源变量提升不明显", "结论不够强", "强调解释性洞察和局限"],
    ]
    add_table_like(slide, rows, 0.75, 1.45, [3.0, 3.5, 5.0], 0.78)
    add_footer(slide, 16)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    set_fill(bg, NAVY)
    add_text(slide, "总结", 0.9, 0.85, 2.0, 0.5, 24, CYAN, True)
    add_text(slide, "我们已经从“项目想法”推进到\n“真实 Lending Club 数据分析”", 0.9, 1.65, 8.8, 1.2, 34, WHITE, True)
    add_metric(slide, 0.95, 3.65, 2.25, 1.0, "2.26M", "真实贷款记录", CYAN)
    add_metric(slide, 3.55, 3.65, 2.25, 1.0, "21.27%", "总体违约率", ORANGE)
    add_metric(slide, 6.15, 3.65, 2.25, 1.0, "A→G", "风险等级梯度", GREEN)
    add_text(slide, "下一阶段：接入 ACS/Fed，解释地区和宏观环境如何影响违约风险。", 0.95, 5.65, 10.8, 0.45, 20, RGBColor(203, 213, 225), True)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    cover(prs)
    slide_motivation(prs)
    slide_questions(prs)
    slide_related(prs)
    slide_data_status(prs)
    slide_pipeline(prs)
    slide_progress(prs)
    slide_finding(prs, 1, "贷款等级与违约率", ["A 等级违约率约 6.57%", "G 等级违约率约 51.57%", "等级体系对违约风险有强区分力"], "lc_default_rate_by_grade.png", BLUE)
    slide_finding(prs, 2, "利率与违约率", ["<8% 利率分箱违约率约 6.38%", ">=24% 利率分箱违约率约 48.56%", "高风险贷款呈现明显违约集中"], "lc_default_rate_by_interest_bin.png", ORANGE)
    slide_finding(prs, 3, "FICO 与违约率", ["FICO 越低，违约率通常越高", "传统信用评分仍有解释力", "多源数据用于补充环境信息"], "lc_default_rate_by_fico_bin.png", GREEN)
    slide_finding(prs, 4, "时间趋势", ["按发放年份观察信贷周期风险", "后续与 Fed 利率、通胀、失业率对齐", "解释宏观环境造成的风险波动"], "lc_default_rate_by_year.png", PURPLE)
    slide_finding(prs, 5, "州级差异", ["州级违约率存在明显差异", "MS 违约率最高，约 28.37%", "后续用 ACS 指标解释地区差异"], "lc_default_rate_top_states.png", CYAN)
    slide_attempts(prs)
    slide_gantt(prs)
    slide_plan(prs)
    slide_risks(prs)
    slide_summary(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
