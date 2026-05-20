#!/usr/bin/env python3
"""
专业风格中期汇报 PPT
现代设计理念：大胆配色、大标题、多样化布局、视觉元素
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

ROOT = Path(__file__).parent.parent
FIGURES = ROOT / "outputs" / "figures"
OUTPUT = ROOT / "slides" / "midterm_professional.pptx"

# 现代金融主题配色 - 深蓝+金红强调
COLORS = {
    "navy_dark": RGBColor(20, 40, 85),      # 深海军蓝
    "navy": RGBColor(30, 64, 135),          # 海军蓝
    "blue": RGBColor(59, 130, 246),         # 亮蓝
    "gold": RGBColor(251, 191, 36),         # 金色
    "coral": RGBColor(249, 97, 103),        # 珊瑚红
    "emerald": RGBColor(16, 185, 129),      # 翡翠绿
    "white": RGBColor(255, 255, 255),
    "light": RGBColor(248, 250, 252),
    "gray_light": RGBColor(148, 163, 184),
    "gray": RGBColor(71, 85, 105),
    "dark": RGBColor(15, 23, 42)
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_shape(slide, shape_type, x, y, w, h, fill=None, line=None):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=None, bold=False, align=None):
    """添加文本"""
    if color is None:
        color = COLORS["gray"]
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = textbox.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align:
        p.alignment = align
    return textbox


def add_numbered_card(slide, number, title, text, x, y, w, h, color):
    """添加带数字的卡片"""
    # 背景
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, COLORS["white"], COLORS["gray_light"])
    # 色条
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.12, h, color)
    # 数字
    add_text(slide, f"{number:02d}", x + 0.18, y + 0.12, 0.5, 0.4, 28, color, True)
    # 标题
    add_text(slide, title, x + 0.75, y + 0.12, w - 0.9, 0.35, 16, COLORS["navy"], True)
    # 内容
    add_text(slide, text, x + 0.75, y + 0.52, w - 0.9, h - 0.65, 12, COLORS["gray"])


def add_stat_card(slide, value, label, x, y, w, h, color):
    """添加统计卡片"""
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, color)
    add_text(slide, value, x + 0.1, y + 0.15, w - 0.2, 0.55, 36, COLORS["white"], True, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, y + 0.75, w - 0.2, 0.3, 11, COLORS["white"], False, PP_ALIGN.CENTER)


def add_image(slide, image_name, x, y, w, h):
    """添加图片"""
    path = FIGURES / image_name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


# ========== 1. 封面页 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
# 渐变背景模拟
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["navy_dark"])
# 装饰圆
add_shape(slide, MSO_SHAPE.OVAL, 9.5, -1.5, 5, 5, COLORS["navy"])
# 标签
add_text(slide, "AI DATA FOUNDATION", 0.7, 0.6, 3, 0.3, 11, COLORS["blue"], True)
# 主标题
add_text(slide, "基于多源数据的", 0.7, 1.3, 12, 0.7, 48, COLORS["white"], True)
add_text(slide, "个人贷款违约风险检测", 0.7, 2.0, 12, 0.7, 48, COLORS["gold"], True)
# 副标题
add_text(slide, "融合 Lending Club · FRED 宏观 · USDA ERS 州级数据", 0.7, 2.9, 8, 0.4, 16, COLORS["gray_light"])
# 统计卡片
add_stat_card(slide, "226万", "贷款记录", 0.7, 4.5, 2.3, 1.2, COLORS["blue"])
add_stat_card(slide, "21.27%", "违约率", 3.2, 4.5, 2.3, 1.2, COLORS["coral"])
add_stat_card(slide, "47", "季度数据", 5.7, 4.5, 2.3, 1.2, COLORS["emerald"])
add_stat_card(slide, "50", "州覆盖", 8.2, 4.5, 2.3, 1.2, COLORS["navy"])
# 日期
add_text(slide, "2026.05.21", 0.7, 6.8, 2, 0.3, 12, COLORS["gray_light"])

# ========== 2. 研究问题 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["light"])
add_text(slide, "研究问题", 0.7, 0.4, 4, 0.5, 14, COLORS["gray_light"], True)
add_text(slide, "违约风险同时受个体、地区与宏观环境三重影响", 0.7, 0.9, 11, 0.7, 36, COLORS["navy_dark"], True)

# 三个层面卡片
add_numbered_card(slide, 1, "个体层面", "贷款等级 A→G\n违约率 6.57%→51.57%", 0.7, 2.0, 3.9, 2.2, COLORS["blue"])
add_numbered_card(slide, 2, "地区层面", "贫困率与违约相关 0.398\n控制利率后仍显著", 4.7, 2.0, 3.9, 2.2, COLORS["emerald"])
add_numbered_card(slide, 3, "宏观层面", "联邦基金利率相关 0.844\n47个季度捕捉周期", 8.7, 2.0, 3.9, 2.2, COLORS["coral"])

# 公式
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.0, 4.8, 9.3, 1.2, COLORS["navy"])
add_text(slide, "违约风险 = 个体特征 + 地区经济 + 宏观周期", 2.2, 5.2, 8.9, 0.6, 24, COLORS["white"], True, PP_ALIGN.CENTER)

# ========== 3. 数据全景 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["light"])
add_text(slide, "数据资产", 0.7, 0.4, 3, 0.5, 14, COLORS["gray_light"], True)
add_text(slide, "三源数据已在个人、州级、季度层面完成融合", 0.7, 0.9, 11, 0.6, 30, COLORS["navy_dark"], True)

# 数据源表格
sources = [
    ["Lending Club", "2,260,701", "151", "个人"],
    ["FRED 宏观", "47", "3", "季度"],
    ["USDA ERS", "50", "3", "州级"]
]

headers = ["数据源", "记录数", "变量数", "粒度"]
y = 1.8
col_widths = [2.5, 2, 1.5, 1.5]
x_start = 1.5

# 表头
for i, h in enumerate(headers):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x_start + sum(col_widths[:i]), y, col_widths[i], 0.5, COLORS["navy"])
    add_text(slide, h, x_start + sum(col_widths[:i]) + 0.05, y + 0.1, col_widths[i] - 0.1, 0.3, 13, COLORS["white"], True, PP_ALIGN.CENTER)

# 数据行
for row_idx, row in enumerate(sources):
    for col_idx, cell in enumerate(row):
        y_pos = y + 0.5 + row_idx * 0.5
        add_text(slide, cell, x_start + sum(col_widths[:col_idx]) + 0.1, y_pos + 0.05, col_widths[col_idx] - 0.2, 0.4, 12, COLORS["gray"])

# 关键指标
add_text(slide, "融合成果", 7.5, 1.8, 1.5, 0.3, 14, COLORS["navy"], True)
add_text(slide, "• 1,367,578 已完结样本", 7.5, 2.2, 5, 0.3, 13, COLORS["gray"])
add_text(slide, "• 50 州完整匹配", 7.5, 2.55, 5, 0.3, 13, COLORS["gray"])
add_text(slide, "• 47 季度时间序列", 7.5, 2.9, 5, 0.3, 13, COLORS["gray"])
add_text(slide, "• 州+季度双重聚合", 7.5, 3.25, 5, 0.3, 13, COLORS["gray"])

# 方法框
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, 4.5, 11.3, 2.2, COLORS["white"], COLORS["gray_light"])
add_text(slide, "融合方法", 1.2, 4.7, 2, 0.3, 14, COLORS["navy"], True)
add_text(slide, "标签定义：Fully Paid(0) / Charged Off+Default(1)", 1.4, 5.1, 5, 0.35, 12, COLORS["gray"])
add_text(slide, "州级聚合：按 addr_state 计算违约率、平均利率、FICO", 1.4, 5.5, 5, 0.35, 12, COLORS["gray"])
add_text(slide, "季度聚合：按 issue_quarter 与 FRED 指标对齐", 1.4, 5.9, 5, 0.35, 12, COLORS["gray"])

# ========== 4. 等级梯度发现 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["light"])
add_text(slide, "发现 01", 0.7, 0.4, 1.5, 0.4, 12, COLORS["blue"], True)
add_text(slide, "贷款等级强区分风险：从 A 等级 6.57% 到 G 等级 51.57%", 0.7, 0.8, 11, 0.6, 32, COLORS["navy_dark"], True)

# 图表
add_image(slide, "lc_default_rate_by_grade.png", 0.7, 1.6, 5.5, 4.0)

# 关键数据卡片
add_stat_card(slide, "6.57%", "A 等级违约率", 6.7, 1.6, 2.5, 1.0, COLORS["emerald"])
add_stat_card(slide, "51.57%", "G 等级违约率", 9.5, 1.6, 2.5, 1.0, COLORS["coral"])

# 解读
add_text(slide, "核心发现", 6.7, 2.8, 2, 0.3, 14, COLORS["navy"], True)
add_text(slide, "• 等级梯度明显：每升一级违约率增加 7-9 个百分点", 6.7, 3.2, 5.3, 0.4, 13, COLORS["gray"])
add_text(slide, "• 平台风险定价有效，等级与利率高度相关", 6.7, 3.65, 5.3, 0.4, 13, COLORS["gray"])
add_text(slide, "• 为多源分析提供重要的基准控制变量", 6.7, 4.1, 5.3, 0.4, 13, COLORS["gray"])

# 来源
add_text(slide, "来源：Lending Club accepted loans (2007-2018), N=1,367,578", 0.7, 5.8, 12, 0.25, 10, COLORS["gray_light"])

# ========== 5. 利率发现 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["light"])
add_text(slide, "发现 02", 0.7, 0.4, 1.5, 0.4, 12, COLORS["blue"], True)
add_text(slide, "利率与违约率高度相关：高风险贷款呈现明显违约集中", 0.7, 0.8, 11, 0.6, 32, COLORS["navy_dark"], True)

add_image(slide, "lc_default_rate_by_interest_bin.png", 0.7, 1.6, 5.5, 4.0)

add_stat_card(slide, "6.38%", "<8% 违约率", 6.7, 1.6, 2.3, 0.9, COLORS["emerald"])
add_stat_card(slide, "48.56%", "≥24% 违约率", 9.2, 1.6, 2.3, 0.9, COLORS["coral"])

add_text(slide, "核心发现", 6.7, 2.7, 2, 0.3, 14, COLORS["navy"], True)
add_text(slide, "• 利率已包含平台风险判断", 6.7, 3.1, 5.3, 0.4, 13, COLORS["gray"])
add_text(slide, "• 高利率贷款仍集中违约", 6.7, 3.55, 5.3, 0.4, 13, COLORS["gray"])
add_text(slide, "• 后续分析需控制利率以评估外部变量增量", 6.7, 4.0, 5.3, 0.4, 13, COLORS["gray"])

add_text(slide, "来源：Lending Club 按利率分箱统计", 0.7, 5.8, 12, 0.25, 10, COLORS["gray_light"])

# ========== 6. 宏观环境发现 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["light"])
add_text(slide, "发现 03", 0.7, 0.4, 1.5, 0.4, 12, COLORS["coral"], True)
add_text(slide, "宏观金融环境与违约率高度相关：联邦基金利率季度相关系数 0.844", 0.7, 0.8, 11, 0.6, 28, COLORS["navy_dark"], True)

add_image(slide, "lc_fred_quarterly_overlay.png", 0.7, 1.6, 5.5, 4.0)

add_stat_card(slide, "0.844", "季度相关系数", 6.7, 1.6, 2.2, 0.9, COLORS["coral"])
add_stat_card(slide, "47", "季度时间点", 9.1, 1.6, 2.2, 0.9, COLORS["blue"])

add_text(slide, "核心发现", 6.7, 2.7, 2, 0.3, 14, COLORS["navy"], True)
add_text(slide, "• 47个季度捕捉完整金融周期", 6.7, 3.1, 5.3, 0.4, 13, COLORS["gray"])
add_text(slide, "• 联邦基金利率与违约率高度相关", 6.7, 3.55, 5.3, 0.4, 13, COLORS["gray"])
add_text(slide, "• 宏观变量值得纳入风险分析框架", 6.7, 4.0, 5.3, 0.4, 13, COLORS["gray"])

add_text(slide, "来源：Lending Club 季度违约率 × FRED 季度宏观指标 (2007-2018)", 0.7, 5.8, 12, 0.25, 10, COLORS["gray_light"])

# ========== 7. 地区经济发现 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["light"])
add_text(slide, "发现 04", 0.7, 0.4, 1.5, 0.4, 12, COLORS["emerald"], True)
add_text(slide, "地区经济变量在控制贷款结构后仍有解释力", 0.7, 0.8, 11, 0.6, 32, COLORS["navy_dark"], True)

add_image(slide, "lc_state_default_residual_interest_vs_poverty.png", 0.7, 1.6, 5.5, 2.3)
add_image(slide, "lc_state_default_residual_fico_vs_poverty.png", 0.7, 4.0, 5.5, 2.3)

add_text(slide, "关键数据", 6.7, 1.6, 2, 0.3, 14, COLORS["navy"], True)
add_text(slide, "原始相关：r = 0.398", 6.7, 2.0, 2.5, 0.35, 13, COLORS["gray"])
add_text(slide, "控制利率后：r = 0.362", 6.7, 2.4, 2.5, 0.35, 13, COLORS["gray"])
add_text(slide, "控制 FICO 后：r = 0.459", 6.7, 2.8, 2.5, 0.35, 13, COLORS["gray"])

add_text(slide, "核心发现", 6.7, 3.4, 2, 0.3, 14, COLORS["navy"], True)
add_text(slide, "• 贫困率与违约率正相关", 6.7, 3.8, 5.3, 0.35, 13, COLORS["gray"])
add_text(slide, "• 控制平均利率后相关性仍显著", 6.7, 4.2, 5.3, 0.35, 13, COLORS["gray"])
add_text(slide, "• 地区经济有独立于个体特征的解释力", 6.7, 4.6, 5.3, 0.35, 13, COLORS["gray"])
add_text(slide, "• 50 个州完整匹配分析", 6.7, 5.0, 5.3, 0.35, 13, COLORS["gray"])

add_text(slide, "来源：Lending Club × USDA ERS (50 州匹配)", 0.7, 6.5, 12, 0.25, 10, COLORS["gray_light"])

# ========== 8. 组合风险发现 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["light"])
add_text(slide, "发现 05", 0.7, 0.4, 1.5, 0.4, 12, COLORS["coral"], True)
add_text(slide, "组合风险分层能精确识别高风险人群", 0.7, 0.8, 11, 0.6, 32, COLORS["navy_dark"], True)

add_image(slide, "lc_top_risk_grade_purpose_segments.png", 0.7, 1.6, 5.5, 4.0)

add_stat_card(slide, "52.99%", "G×debt_consolidation", 6.7, 1.6, 2.3, 0.9, COLORS["coral"])
add_stat_card(slide, "51.64%", "G×60months", 9.1, 1.6, 2.3, 0.9, COLORS["coral"])

add_text(slide, "核心发现", 6.7, 2.7, 2, 0.3, 14, COLORS["navy"], True)
add_text(slide, "• 组合维度比单一变量更精确", 6.7, 3.1, 5.3, 0.35, 13, COLORS["gray"])
add_text(slide, "• G×debt_consolidation 违约率 52.99%", 6.7, 3.5, 5.3, 0.35, 13, COLORS["gray"])
add_text(slide, "• 可直接用于精准识别和定价", 6.7, 3.9, 5.3, 0.35, 13, COLORS["gray"])

add_text(slide, "来源：Lending Club 组合风险分层分析", 0.7, 5.8, 12, 0.25, 10, COLORS["gray_light"])

# ========== 9. 方法论价值 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["light"])
add_text(slide, "方法论价值", 0.7, 0.4, 2, 0.5, 14, COLORS["gray_light"], True)
add_text(slide, "多源融合证明了什么", 0.7, 0.9, 11, 0.6, 36, COLORS["navy_dark"], True)

value_cards = [
    ("✓", "个人变量有效", "等级梯度明显，是基准控制变量", COLORS["blue"]),
    ("✓", "地区变量独立", "控制贷款结构后仍显著", COLORS["emerald"]),
    ("✓", "宏观变量相关", "联邦基金利率相关 0.844", COLORS["coral"]),
    ("✓", "组合更精确", "可识别高风险人群", COLORS["gold"])
]

y = 1.8
for icon, title, text, color in value_cards:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, y, 11.3, 1.15, COLORS["white"], COLORS["gray_light"])
    add_shape(slide, MSO_SHAPE.OVAL, 1.2, y + 0.2, 0.7, 0.7, color)
    add_text(slide, icon, 1.25, y + 0.25, 0.6, 0.6, 28, COLORS["white"], True, PP_ALIGN.CENTER)
    add_text(slide, title, 2.1, y + 0.15, 2.5, 0.35, 16, COLORS["navy"], True)
    add_text(slide, text, 2.1, y + 0.55, 9, 0.4, 13, COLORS["gray"])
    y += 1.25

# ========== 10. 结论页 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, COLORS["navy_dark"])
add_shape(slide, MSO_SHAPE.OVAL, 9, -1.5, 5, 5, COLORS["navy"])

add_text(slide, "结论", 0.7, 0.6, 2, 0.5, 20, COLORS["blue"], True)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0.7, 1.15, 2, 0.04, COLORS["blue"])

conclusions = [
    ("01", "贷款等级是强风险分层工具", "A 等级违约率 6.57%，G 等级 51.57%"),
    ("02", "宏观环境与违约率高度相关", "联邦基金利率季度相关 0.844"),
    ("03", "地区经济有独立解释力", "控制贷款结构后贫困率仍相关 0.362"),
    ("04", "多源融合能识别高风险人群", "G×debt_consolidation 违约率 52.99%")
]

y = 1.5
for num, title, text in conclusions:
    add_text(slide, num, 0.7, y, 0.6, 0.4, 24, COLORS["blue"], True)
    add_text(slide, title, 1.5, y, 4, 0.35, 16, COLORS["white"], True)
    add_text(slide, text, 1.5, y + 0.4, 8, 0.3, 13, COLORS["gray_light"])
    y += 0.8

add_text(slide, "Questions & Answers", 0.7, 5.2, 4, 0.5, 24, COLORS["gold"], True)
add_text(slide, "项目代码 | 数据与图表 | outputs/ 目录", 0.7, 5.8, 6, 0.3, 12, COLORS["gray_light"])

# 保存
prs.save(OUTPUT)
print(f"已生成专业风格 PPT: {OUTPUT}")
