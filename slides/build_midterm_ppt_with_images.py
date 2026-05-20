#!/usr/bin/env python3
"""
中期汇报 PPT 生成 - 详细版（带真实图表）
包含方法论、技术细节、完整发现、真实图表
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches
import os

# 配色方案
COLORS = {
    'navy': RGBColor(30, 39, 97),
    'blue': RGBColor(46, 117, 182),
    'coral': RGBColor(249, 97, 103),
    'emerald': RGBColor(16, 185, 129),
    'gold': RGBColor(249, 231, 149),
    'dark_gray': RGBColor(52, 52, 52),
    'light_gray': RGBColor(128, 128, 128),
    'white': RGBColor(255, 255, 255),
}

# 图表路径
FIGURES_PATH = '/Users/bytedance/Desktop/DB/outputs/figures'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['navy']
    background.line.color.rgb = COLORS['navy']

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(42)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.text = subtitle
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = COLORS['gold']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    info_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(0.5))
    tf = info_box.text_frame
    tf.text = "AI Data Foundation Project | 2026-05-21"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = COLORS['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_section_divider(prs, number, title, subtitle=""):
    """章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['navy']
    background.line.color.rgb = COLORS['navy']

    number_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(2.5), Inches(2))
    tf = number_box.text_frame
    tf.text = f"{number:02d}"
    tf.paragraphs[0].font.size = Pt(100)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(100, 110, 150)

    title_box = slide.shapes.add_textbox(Inches(3.5), Inches(3), Inches(8.83), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['white']

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(3.5), Inches(4), Inches(8.83), Inches(0.6))
        tf = subtitle_box.text_frame
        tf.text = subtitle
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = COLORS['gold']

    return slide

def add_content_slide(prs, title, content_items):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['navy']
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['navy']

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.7))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = item
        p.font.size = Pt(17)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_after = Pt(10)

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """双列内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['navy']
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['navy']

    left_header = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(0.05), Inches(0.4))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = COLORS['blue']

    left_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    tf.text = left_title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['blue']

    left_content = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.8), Inches(5))
    tf = left_content.text_frame
    for i, item in enumerate(left_items):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_after = Pt(8)

    right_header = slide.shapes.add_shape(1, Inches(6.8), Inches(1.2), Inches(0.05), Inches(0.4))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = COLORS['coral']

    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.2), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    tf.text = right_title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['coral']

    right_content = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    tf = right_content.text_frame
    for i, item in enumerate(right_items):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_after = Pt(8)

    return slide

def add_table_slide(prs, title, table_data):
    """表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['navy']
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['navy']

    rows = len(table_data)
    cols = len(table_data[0])
    table_width = Inches(12)
    table_height = Inches(5)
    x_start = Inches(0.67)
    y_start = Inches(1.3)

    table = slide.shapes.add_table(rows, cols, x_start, y_start, table_width, table_height).table

    for r, row_data in enumerate(table_data):
        for c, cell_data in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].font.size = Pt(13) if r > 0 else Pt(15)
            cell.text_frame.paragraphs[0].font.bold = (r == 0)
            cell.text_frame.paragraphs[0].font.color.rgb = COLORS['white'] if r == 0 else COLORS['dark_gray']
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c in [0, 1, 3] else PP_ALIGN.LEFT

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['navy']
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 249, 250)

    return slide

def add_chart_slide(prs, title, image_filename, key_findings, chart_side="left"):
    """图表页 - 带真实图片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['navy']
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['navy']

    # 插入真实图表
    image_path = os.path.join(FIGURES_PATH, image_filename)
    if os.path.exists(image_path):
        if chart_side == "left":
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.3), width=Inches(6.5))
            findings_x = Inches(7.2)
        else:
            slide.shapes.add_picture(image_path, Inches(6.8), Inches(1.3), width=Inches(6))
            findings_x = Inches(0.7)
    else:
        # 如果图片不存在，显示占位符
        placeholder = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(6.5), Inches(4))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(248, 249, 250)
        placeholder.line.color.rgb = COLORS['light_gray']
        findings_x = Inches(7.2)

    # 关键发现
    findings_box = slide.shapes.add_textbox(findings_x, Inches(1.3), Inches(5.5), Inches(5.5))
    tf = findings_box.text_frame
    tf.word_wrap = True

    for i, finding in enumerate(key_findings):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = finding
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_after = Pt(12)

    return slide

# ==================== 开始创建 PPT ====================

# 1. 封面
add_title_slide(prs,
    "基于多源数据的个人贷款违约风险检测",
    "Multi-Source Data Fusion for Loan Default Risk Detection\n中期汇报 | 2026-05-21"
)

# 2. 章节分隔
add_section_divider(prs, 1, "研究动机与问题定义")

# 3. 核心问题
add_content_slide(prs,
    "核心研究问题",
    [
        "🤔 同样收入、同样信用评分的两个人，违约风险会一样吗？",
        "",
        "• 一个住在密西西比州（贫困率19%），一个住在加州（贫困率13%）",
        "• 一个在2008年低利率时期贷款，一个在2006年高利率时期贷款",
        "",
        "📊 传统模型：违约风险 = f(个人特征)",
        "🎯 本研究：违约风险 = f(个人特征 + 地区经济 + 宏观周期)",
    ]
)

# 4. 研究动机
add_two_column_slide(prs,
    "为什么这项研究重要？",
    "📊 实际价值",
    [
        "• 个人贷款违约直接影响金融机构坏账损失",
        "• 影响借款人的信用定价和融资成本",
        "• 风险定价不准会导致不公平现象",
        "• 多源数据可以提升定价精准度",
    ],
    "🎯 科学问题",
    [
        "• 地区经济如何影响个人违约风险？",
        "• 宏观周期与个人风险有何关系？",
        "• 各层变量的独立解释力有多大？",
        "• 多源融合能带来多少增量价值？",
    ]
)

# 5. 核心挑战
add_content_slide(prs,
    "项目面临的核心挑战",
    [
        "📁 挑战一：数据规模大",
        "   Lending Club: 226万条记录，151个字段，文件2.8GB",
        "",
        "🏷️ 挑战二：标签复杂",
        "   贷款状态有10多种：Fully Paid, Charged Off, Current, Late...",
        "   未完结状态无法判断最终是否违约",
        "",
        "🔗 挑战三：多源数据粒度不一致",
        "   贷款数据(个人级) ← 州级数据(50州) ← 宏观数据(时间序列)",
    ]
)

# 6. 章节分隔
add_section_divider(prs, 2, "相关工作与定位")

# 7. 方法对比
add_table_slide(prs, "现有方法对比",
    [
        ["方法类型", "代表技术", "优点", "局限"],
        ["传统信用评分", "FICO、收入、负债率", "解释性强", "只看个人"],
        ["机器学习", "LR、RF、XGBoost", "预测能力强", "容易黑盒化"],
        ["多源分析", "地区收入、失业率", "视角全面", "数据对齐难"],
        ["本研究", "个体+地区+宏观", "先做可解释分析", "再考虑预测"],
    ]
)

# 8. 章节分隔
add_section_divider(prs, 3, "数据全景与融合")

# 9. 数据源表格
add_table_slide(prs, "四源数据融合",
    [
        ["数据源", "记录数", "变量数", "作用"],
        ["Lending Club", "226万", "151", "个体贷款记录、违约标签"],
        ["FRED 宏观", "47季度", "3", "联邦基金利率、CPI、失业率"],
        ["USDA ERS", "50州", "3", "州级收入、贫困率、失业率"],
        ["Home Credit", "30万", "122+", "补充信贷场景对照"],
    ]
)

# 10. 融合成果
add_content_slide(prs,
    "数据融合成果",
    [
        "✅ 已完结样本筛选",
        "   从 2,260,701 条原始记录 → 1,367,578 条已完结样本",
        "   排除 Current、Late 等未完结状态，确保标签质量",
        "",
        "✅ 时空完整覆盖",
        "   • 50 个州：完整匹配州级经济数据",
        "   • 47 个季度：2007-Q2 至 2018-Q4",
        "",
        "✅ 多源对齐完成",
        "   个人级 ←→ 州级数据（通过 addr_state）",
        "   个人级 ←→ 宏观数据（通过 issue_quarter）",
    ]
)

# 11. 章节分隔 - 方法论
add_section_divider(prs, 4, "方法论与技术实现", "Methodology & Implementation")

# 12. 技术架构
add_content_slide(prs,
    "五层技术架构",
    [
        "┌─────────────────────────────────────────────────────────────┐",
        "│ Layer 1: 数据获取层                                         │",
        "│ Kaggle API + FRED API + USDA ERS 爬取                        │",
        "└─────────────────────────────────────────────────────────────┘",
        "                              ↓",
        "┌─────────────────────────────────────────────────────────────┐",
        "│ Layer 2: 数据清洗层   流式CSV | 状态过滤 | 缺失值处理        │",
        "└─────────────────────────────────────────────────────────────┘",
        "                              ↓",
        "┌─────────────────────────────────────────────────────────────┐",
        "│ Layer 3: 特征工程层   标签构造 | 时间提取 | 州级聚合         │",
        "└─────────────────────────────────────────────────────────────┘",
        "                              ↓",
        "┌─────────────────────────────────────────────────────────────┐",
        "│ Layer 4: 多源融合层   州匹配 | 季度对齐 | 残差计算           │",
        "└─────────────────────────────────────────────────────────────┘",
        "                              ↓",
        "┌─────────────────────────────────────────────────────────────┐",
        "│ Layer 5: 统计分析层   分组聚合 | 相关系数 | 交叉分析         │",
        "└─────────────────────────────────────────────────────────────┘",
    ]
)

# 13. 方法1：流式处理
add_two_column_slide(prs,
    "方法一：处理大规模 CSV",
    "❌ 问题",
    [
        "• 文件 2.8GB，内存溢出",
        "• pd.read_csv() 直接读取失败",
    ],
    "✅ 流式处理方案",
    [
        "for chunk in pd.read_csv(",
        "    path, chunksize=10000):",
        "    processed = process(",
        "        chunk)",
        "    accumulate(",
        "        processed)",
        "",
        "每次只读1万条，",
        "处理完立即释放",
    ]
)

# 14. 方法2：违约标签
add_content_slide(prs,
    "方法二：违约标签构造（关键决策）",
    [
        "🔍 原始数据有10+种状态，需要准确定义违约",
        "",
        "✅ 已完结状态（保留）:",
        "   • Fully Paid → 0（不违约）",
        "   • Charged Off / Default → 1（违约）",
        "",
        "❌ 未完结状态（排除）:",
        "   • Current、Late、In Payment Plan",
        "   → 无法判断最终结果，会导致偏差",
        "",
        "📊 结果：从 226万 → 136万已完结样本",
        "   样本量减少，但标签质量提高",
    ]
)

# 15. 方法3：时间对齐
add_two_column_slide(prs,
    "方法三：多源时间对齐",
    "📅 贷款数据时间格式",
    [
        "原始: 'Dec-2017'",
        "解析: pd.to_datetime(",
        "  issue_d, format='%b-%Y')",
        "提取: year, quarter",
        "",
        "覆盖: 2007-Q2 ~ 2018-Q4",
        "共 47 个季度",
    ],
    "🔄 宏观数据对齐",
    [
        "FRED 是月度数据",
        "需要转成季度",
        "",
        "方法: 取季度最后",
        "一个月的值",
        "",
        "代表发放时的",
        "宏观环境",
    ]
)

# 16. 方法4：控制变量
add_content_slide(prs,
    "方法四：控制变量分析（核心创新）",
    [
        "🤔 问题：贫困州违约率高，是因为那里的人信用分低吗？",
        "",
        "💡 解决思路：计算「残差违约率」——剔除其他因素后的剩余变异",
        "",
        "📐 分析步骤：",
        "   Step 1: 建立利率→违约率的 Logistic 回归模型",
        "   Step 2: 预测每个州的「预期违约率」",
        "   Step 3: 计算残差 = 实际违约率 - 预期违约率",
        "   Step 4: 分析残差与贫困率的相关性",
        "",
        "📊 结果：控制利率后，残差与贫困率仍相关 r=0.362",
        "   ✅ 地区经济有独立于个人特征的解释力！",
    ]
)

# 17. 章节分隔
add_section_divider(prs, 5, "核心发现", "Key Findings")

# 18. 发现1：等级梯度
add_chart_slide(prs,
    "发现一：贷款等级强区分违约风险",
    "lc_default_rate_by_grade.png",
    [
        "📈 关键数据：",
        "   A 等级：6.57%",
        "   G 等级：51.57%",
        "   相差 8 倍！",
        "",
        "💡 意义：",
        "   • 平台风险定价有效",
        "   • 个人特征仍是强因子",
        "   • 等级作为控制变量",
    ],
    chart_side="left"
)

# 19. 发现2：利率相关
add_chart_slide(prs,
    "发现二：利率与违约率高度相关",
    "lc_default_rate_by_interest_bin.png",
    [
        "📈 关键数据：",
        "   低利率（<8%）：6.38%",
        "   高利率（≥24%）：48.56%",
        "",
        "💡 意义：",
        "   • 平台定价包含风险判断",
        "   • 高利率贷款违约集中",
        "   • 需控制利率做后续分析",
    ],
    chart_side="left"
)

# 20. 发现3：宏观环境
add_chart_slide(prs,
    "发现三：宏观金融环境与违约率高度相关",
    "lc_fred_quarterly_overlay.png",
    [
        "📈 关键数据：",
        "   季度相关系数：0.844",
        "   这是非常高的相关性",
        "",
        "⚠️ 诚实说明：",
        "   • 这是相关性，非因果",
        "   • 可能是共同因素导致",
        "",
        "💡 意义：",
        "   宏观环境值得纳入框架",
    ],
    chart_side="left"
)

# 21. 发现4：地区经济
add_chart_slide(prs,
    "发现四：地区经济有独立解释力",
    "lc_state_default_residual_interest_vs_poverty.png",
    [
        "📈 控制变量结果：",
        "   原始：r = 0.398",
        "   控制利率：r = 0.362",
        "   控制FICO：r = 0.459",
        "",
        "🎯 核心发现：",
        "   即使两个人信用分相同、",
        "   利率相同，生活在贫困率",
        "   更高的州，违约风险仍高",
        "",
        "   地区经济有独立解释力！",
    ],
    chart_side="left"
)

# 22. 发现5：组合风险
add_chart_slide(prs,
    "发现五：组合风险分层识别高风险人群",
    "lc_top_risk_grade_purpose_segments.png",
    [
        "🎯 最高风险组合：",
        "   G × debt_consolidation",
        "   违约率：52.99%",
        "",
        "💡 实际意义：",
        "   • 每两笔就有一笔违约",
        "   • 可直接用于风控",
        "   • 组合 > 单一变量",
    ],
    chart_side="left"
)

# 23. 章节分隔
add_section_divider(prs, 6, "项目进度与后续计划")

# 24. 项目进度
add_content_slide(prs,
    "项目时间线与当前进度",
    [
        "✅ 已完成（4/20 - 5/20）：",
        "   • 开题与问题定义",
        "   • 项目仓库搭建",
        "   • Lending Club 数据分析（226万条）",
        "   • FRED/ERS 数据融合",
        "   • 风险分层分析",
        "   • Home Credit 数据下载",
        "",
        "🔄 进行中（5/21 - 5/27）：",
        "   • 控制变量建模",
        "",
        "⏳ 计划中（5/27 - 6/17）：",
        "   • 最终可视化与结论提炼",
        "   • PPT/报告/视频打磨",
        "",
        "📍 当前节点：2026年5月20日",
    ]
)

# 25. 风险管理
add_table_slide(prs, "风险与备选方案",
    [
        ["风险", "影响", "备选方案", "状态"],
        ["Home Credit授权", "数据缺失", "已有LC+FRED+ERS", "✅解决"],
        ["相关性≠因果", "解释有限", "控制变量模型", "🔄进行"],
        ["时间紧张", "建模简化", "保多源分析", "📋备用"],
        ["变量提升不显著", "结论弱", "强调方法论", "📋备用"],
    ]
)

# 26. 总结
add_content_slide(prs,
    "总结：核心发现",
    [
        "✅ 发现一：贷款等级是强风险分层工具",
        "   A→G 相差 8 倍",
        "",
        "✅ 发现二：宏观环境与违约率高度相关",
        "   相关系数 0.844",
        "",
        "✅ 发现三：地区经济有独立解释力",
        "   控制结构后仍相关 0.362",
        "",
        "✅ 发现四：多源融合精确识别高风险",
        "   G×debt_consolidation 违约率 52.99%",
        "",
        "🎯 核心贡献：完成个人-地区-宏观三层分析",
        "   证明违约风险不只是个人问题",
    ]
)

# 27. Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = COLORS['navy']
background.line.color.rgb = COLORS['navy']

qa_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(1.5))
tf = qa_box.text_frame
tf.text = "Questions & Answers"
tf.paragraphs[0].font.size = Pt(48)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = COLORS['white']
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.33), Inches(0.8))
tf = contact_box.text_frame
tf.text = "感谢聆听 | Thank You\nGitHub: github.com/Epiphanythu/ai-data-foundation-project"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = COLORS['gold']
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 保存
output_path = '/Users/bytedance/Desktop/DB/slides/midterm_with_charts.pptx'
prs.save(output_path)
print(f"✅ PPT 已生成: {output_path}")
print(f"   总页数: {len(prs.slides)}")
print(f"   包含真实图表: 5张核心发现图表")
print(f"   包含方法论详细讲解: 5页")
